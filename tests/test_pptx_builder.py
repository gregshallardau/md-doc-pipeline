"""Tests for the pptx (PowerPoint) slide builder."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from md_doc.builders.pptx import build


@pytest.fixture()
def repo(tmp_path):
    (tmp_path / ".git").mkdir()
    return tmp_path


def _png(path: Path) -> None:
    from PIL import Image

    Image.new("RGB", (200, 120), "teal").save(path)


def _build(repo: Path, body: str, config: dict) -> Presentation:
    from md_doc.config import load_config

    doc = repo / "talk.md"
    doc.write_text(body, encoding="utf-8")
    out = repo / "talk.pptx"
    cfg = load_config(doc, repo_root=repo)  # merge frontmatter like the real pipeline
    cfg.update(config)
    build(body, cfg, out, doc_path=doc, repo_root=repo)
    return Presentation(str(out))


def _titles(prs: Presentation) -> list[str]:
    return [(s.shapes.title.text if s.shapes.title else "") for s in prs.slides]


def _pics(slide) -> int:
    return sum(1 for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)


def _tables(slide) -> int:
    return sum(1 for sh in slide.shapes if sh.has_table)


DECK = """---
title: Quarterly Review
author: Jane Doe
product: Acme
date: July 2026
---

# Quarterly Review

## Highlights

- Top level
    - Nested one

<!-- notes: speaker note here -->

## Metrics

| M | Q1 | Q2 |
|---|----|----|
| U | 10 | 14 |

## Shot

![dash](pic.png)
"""


def test_title_slide_and_h2_segmentation(repo):
    prs = _build(repo, DECK, {})
    titles = _titles(prs)
    # title slide + 3 H2 content slides
    assert titles[0] == "Quarterly Review"
    assert "Highlights" in titles
    assert "Metrics" in titles
    assert "Shot" in titles


def test_title_slide_has_metadata(repo):
    prs = _build(repo, DECK, {})
    # subtitle placeholder carries product/author/date
    text = "\n".join(
        ph.text_frame.text for ph in prs.slides[0].placeholders if ph.placeholder_format.idx != 0
    )
    assert "Acme" in text and "Jane Doe" in text and "July 2026" in text


def test_bullet_levels(repo):
    prs = _build(repo, DECK, {})
    highlights = next(
        s for s in prs.slides if s.shapes.title and s.shapes.title.text == "Highlights"
    )
    levels = {
        p.level
        for sh in highlights.shapes
        if sh.has_text_frame
        for p in sh.text_frame.paragraphs
        if p.text
    }
    assert 0 in levels and 1 in levels  # nested bullet produced a level-1 paragraph


def test_table_and_image_and_notes(repo):
    _png(repo / "pic.png")
    prs = _build(repo, DECK, {})
    assert any(_tables(s) for s in prs.slides)
    assert any(_pics(s) for s in prs.slides)
    notes = [s.notes_slide.notes_text_frame.text for s in prs.slides if s.has_notes_slide]
    assert any("speaker note here" in n for n in notes)


def test_explicit_slide_break(repo):
    body = "---\ntitle: T\n---\n\n## One\n\nalpha\n\n<!-- slide -->\n\nbeta\n"
    prs = _build(repo, body, {})
    # title slide + One + the explicit-break slide
    assert len(prs.slides._sldIdLst) == 3


def test_slide_split_marker_only(repo):
    body = "---\ntitle: T\n---\n\n## One\n\n## Two\n\n<!-- slide -->\n\nlast\n"
    prs = _build(repo, body, {"slide_split": "marker"})
    # H2 headings do NOT split in marker mode: title slide + one content + marker slide
    assert len(prs.slides._sldIdLst) == 3


def test_slide_size_widescreen(repo):
    prs = _build(repo, DECK, {"slide_size": "16:9"})
    from pptx.util import Inches

    assert prs.slide_width == Inches(13.333)


def test_picks_up_css_theme_and_yaml(repo):
    # YAML frontmatter (title/author/product) and the CSS theme cascade
    # (h1 colour + body font) both flow into the deck.
    (repo / "_pdf-theme.css").write_text(
        ':root { --primary: #CC0066; }\nbody { font-family: "Georgia"; }\nh1 { color: #CC0066; }\n',
        encoding="utf-8",
    )
    body = (
        "---\ntitle: Themed Deck\nauthor: Ada\nproduct: Engine\n---\n\n"
        "# Themed Deck\n\n## Point\n\n- hello\n"
    )
    prs = _build(repo, body, {})
    # YAML → title slide
    assert prs.slides[0].shapes.title.text == "Themed Deck"
    sub = "\n".join(
        ph.text_frame.text for ph in prs.slides[0].placeholders if ph.placeholder_format.idx != 0
    )
    assert "Ada" in sub and "Engine" in sub
    # CSS → content-slide title colour
    content = next(s for s in prs.slides if s.shapes.title and s.shapes.title.text == "Point")
    from pptx.dml.color import RGBColor

    assert content.shapes.title.text_frame.paragraphs[0].font.color.rgb == RGBColor(
        0xCC, 0x00, 0x66
    )
    # CSS → body font
    fonts = {
        r.font.name
        for sh in content.shapes
        if sh.has_text_frame
        for p in sh.text_frame.paragraphs
        for r in p.runs
    }
    assert "Georgia" in fonts


def test_mermaid_diagram_embeds_as_picture(repo):
    pytest.importorskip("cairosvg")
    body = '---\ntitle: T\n---\n\n## Chart\n\n```mermaid\npie\n"A" : 60\n"B" : 40\n```\n'
    prs = _build(repo, body, {})
    assert any(_pics(s) for s in prs.slides)
