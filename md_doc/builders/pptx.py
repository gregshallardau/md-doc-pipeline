"""python-pptx builder for .pptx (PowerPoint) output.

Converts rendered Markdown into a slide deck. Unlike the flowing pdf/docx
builders, this one *segments* the document into slides:

- the first ``# H1`` (or config ``title``) becomes a **title slide**;
- each ``# H1`` after that becomes a **section-divider slide**;
- each ``## H2`` becomes a **content slide** (heading → slide title);
- a ``<!-- slide -->`` marker forces a new slide anywhere.

The ``slide_split`` config key selects the strategy: ``h2`` (default) splits on
H2, ``h1`` only on H1, ``marker`` only on ``<!-- slide -->``.

Content within a slide (lists, paragraphs, tables, images, Mermaid diagrams,
code blocks) is laid out top-to-bottom in the body area. A ``<!-- notes: … -->``
comment attaches speaker notes to the current slide.

Public API
----------
    build(rendered_md, config, out_path, *, doc_path=None, repo_root=None)
"""

from __future__ import annotations

import logging
import re
from dataclasses import dataclass, field
from html.parser import HTMLParser
from io import BytesIO
from pathlib import Path
from typing import Any

import markdown
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

from ..docx_theme import _hex_to_rgb, resolve_docx_theme
from ._assets import _EMU_PER_PX, _MERMAID_IMG_RE, _render_mermaid_to_images, _resolve_asset

logger = logging.getLogger(__name__)

_MD_EXTENSIONS = [
    "tables",
    "fenced_code",
    "footnotes",
    "def_list",
    "abbr",
    "attr_list",
    "md_in_html",
]

_SLIDE_BREAK_RE = re.compile(r"<!--\s*slide\s*-->", re.IGNORECASE)
_NOTES_RE = re.compile(r"<!--\s*notes:\s*(.*?)\s*-->", re.IGNORECASE | re.DOTALL)

_SLIDE_SIZES = {"16:9": (Inches(13.333), Inches(7.5)), "4:3": (Inches(10), Inches(7.5))}


# ---------------------------------------------------------------------------
# Slide model
# ---------------------------------------------------------------------------


@dataclass
class _Run:
    text: str
    bold: bool = False
    italic: bool = False
    code: bool = False


@dataclass
class _Para:
    runs: list[_Run] = field(default_factory=list)
    level: int = 0
    kind: str = "para"  # para | bullet | number | code


@dataclass
class _Slide:
    title: str | None = None
    kind: str = "content"  # title | section | content
    paras: list[_Para] = field(default_factory=list)
    images: list[int] = field(default_factory=list)  # mermaid image indices
    files: list[Path] = field(default_factory=list)  # resolved image files
    tables: list[list[list[str]]] = field(default_factory=list)
    notes: list[str] = field(default_factory=list)

    def is_empty(self) -> bool:
        return not (self.paras or self.images or self.files or self.tables or self.title)


# ---------------------------------------------------------------------------
# HTML → slides
# ---------------------------------------------------------------------------


class _SlideParser(HTMLParser):
    """Walk the document HTML and produce a list of :class:`_Slide`."""

    def __init__(self, split: str, doc_path: Path | None, repo_root: Path | None) -> None:
        super().__init__()
        self.convert_charrefs = True
        self._split = split
        self._doc_path = doc_path
        self._repo_root = repo_root
        self.slides: list[_Slide] = []
        self._cur: _Slide | None = None

        self._para: _Para | None = None
        self._bold = False
        self._italic = False
        self._code = False
        self._list_stack: list[str] = []
        self._in_pre = False
        self._pre_text = ""
        self._pending_heading: str | None = None
        self._heading_text = ""

        # table state
        self._in_table = False
        self._cur_row: list[str] | None = None
        self._cur_table: list[list[str]] | None = None
        self._cell = ""
        self._in_cell = False

    # -- slide/paragraph management ------------------------------------------

    def _new_slide(self, title: str | None, kind: str) -> None:
        self._cur = _Slide(title=title, kind=kind)
        self.slides.append(self._cur)
        self._para = None

    def _slide(self) -> _Slide:
        if self._cur is None:
            self._new_slide(None, "content")
        assert self._cur is not None
        return self._cur

    def _start_para(self, kind: str) -> None:
        # Ensure the slide exists *first* — _slide() may create one via
        # _new_slide(), which resets self._para.
        slide = self._slide()
        self._para = _Para(level=max(len(self._list_stack) - 1, 0), kind=kind)
        slide.paras.append(self._para)

    def _add_text(self, text: str) -> None:
        if self._para is None or not text:
            return
        self._para.runs.append(
            _Run(text=text, bold=self._bold, italic=self._italic, code=self._code)
        )

    # -- HTMLParser callbacks ------------------------------------------------

    def handle_starttag(self, tag: str, attrs: list) -> None:
        tag = tag.lower()
        if self._in_cell and tag not in ("td", "th", "tr", "table"):
            return  # inline cell markup collected as text

        if tag in ("h1", "h2", "h3", "h4"):
            self._pending_heading = tag  # capture text, act on end tag
            self._heading_text = ""
        elif tag == "p":
            if not self._in_table:
                self._start_para("para")
        elif tag in ("ul", "ol"):
            self._list_stack.append(tag)
        elif tag == "li":
            kind = "number" if self._list_stack and self._list_stack[-1] == "ol" else "bullet"
            self._start_para(kind)
        elif tag in ("strong", "b"):
            self._bold = True
        elif tag in ("em", "i"):
            self._italic = True
        elif tag == "code" and not self._in_pre:
            self._code = True
        elif tag == "pre":
            self._in_pre = True
            self._pre_text = ""
        elif tag == "img":
            self._embed_image(dict(attrs))
        elif tag == "table":
            self._in_table = True
            self._cur_table = []
        elif tag == "tr" and self._in_table:
            self._cur_row = []
        elif tag in ("td", "th") and self._in_table:
            self._in_cell = True
            self._cell = ""
        elif tag == "br" and self._para is not None:
            self._para.runs.append(_Run(text="\n"))

    def handle_startendtag(self, tag: str, attrs: list) -> None:
        if tag.lower() == "img":
            self._embed_image(dict(attrs))
        else:
            self.handle_starttag(tag, attrs)
            self.handle_endtag(tag)

    def handle_endtag(self, tag: str) -> None:
        tag = tag.lower()
        if tag in ("h1", "h2", "h3", "h4"):
            self._finish_heading(tag)
        elif tag in ("ul", "ol"):
            if self._list_stack:
                self._list_stack.pop()
        elif tag in ("li", "p"):
            self._para = None
        elif tag in ("strong", "b"):
            self._bold = False
        elif tag in ("em", "i"):
            self._italic = False
        elif tag == "code" and not self._in_pre:
            self._code = False
        elif tag == "pre":
            self._in_pre = False
            para = _Para(kind="code")
            para.runs.append(_Run(text=self._pre_text.rstrip("\n"), code=True))
            self._slide().paras.append(para)
            self._pre_text = ""
        elif tag in ("td", "th") and self._in_table:
            self._in_cell = False
            if self._cur_row is not None:
                self._cur_row.append(self._cell.strip())
        elif tag == "tr" and self._in_table:
            if self._cur_row and self._cur_table is not None:
                self._cur_table.append(self._cur_row)
            self._cur_row = None
        elif tag == "table":
            self._in_table = False
            if self._cur_table:
                self._slide().tables.append(self._cur_table)
            self._cur_table = None

    def handle_data(self, data: str) -> None:
        if getattr(self, "_pending_heading", None):
            self._heading_text += data
            return
        if self._in_pre:
            self._pre_text += data
            return
        if self._in_cell:
            self._cell += data
            return
        self._add_text(data)

    def handle_comment(self, data: str) -> None:
        if data.strip() == "slidebreak":
            self._new_slide(None, "content")
            return
        m = _NOTES_RE.match(f"<!--{data}-->")
        if m:
            self._slide().notes.append(m.group(1).strip())

    # -- helpers -------------------------------------------------------------

    def _finish_heading(self, tag: str) -> None:
        text = getattr(self, "_heading_text", "").strip()
        self._pending_heading = None
        if tag == "h1":
            self._new_slide(text, "section")
        elif tag == "h2" and self._split in ("h2",):
            self._new_slide(text, "content")
        else:
            # heading kept inline as a bold lead-in paragraph
            self._start_para("para")
            self._para.runs.append(_Run(text=text, bold=True))
            self._para = None

    def _embed_image(self, attrs: dict) -> None:
        src = attrs.get("src") or ""
        m = _MERMAID_IMG_RE.fullmatch(src)
        if m:
            self._slide().images.append(int(m.group(1)))
            return
        path = _resolve_asset(src, self._doc_path, self._repo_root)
        if path is not None:
            self._slide().files.append(path)
        else:
            logger.warning("pptx: could not resolve image %r — skipped.", src)


# ---------------------------------------------------------------------------
# Rendering slides → python-pptx
# ---------------------------------------------------------------------------


def _layout(prs: Presentation, name: str, fallback: int):
    for lyt in prs.slide_layouts:
        if lyt.name.strip().lower() == name.lower():
            return lyt
    idx = min(fallback, len(prs.slide_layouts) - 1)
    return prs.slide_layouts[idx]


def _rgb(hex_str: str | None) -> RGBColor | None:
    if not hex_str:
        return None
    try:
        r, g, b = _hex_to_rgb(hex_str)
        return RGBColor(r, g, b)
    except Exception:
        return None


def _write_para(tf, para: _Para, first: bool, theme: dict, body_font: str | None) -> None:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.level = min(para.level, 4)
    code_font = theme.get("font_code", "Courier New")
    for run in para.runs:
        r = p.add_run()
        r.text = run.text
        if run.bold:
            r.font.bold = True
        if run.italic:
            r.font.italic = True
        if run.code:
            r.font.name = code_font
        elif body_font:
            r.font.name = body_font
    if para.kind == "code":
        p.font.name = code_font
        p.font.size = Pt(float(theme.get("font_size_code", 12.0)))


def _fill_text_frame(tf, paras: list[_Para], theme: dict, body_font: str | None) -> None:
    tf.word_wrap = True
    try:
        from pptx.enum.text import MSO_AUTO_SIZE

        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    first = True
    for para in paras:
        _write_para(tf, para, first, theme, body_font)
        first = False


def _add_table(slide, rows: list[list[str]], left, top, width, theme: dict) -> Emu:
    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    height = Emu(int(Pt(22).emu * n_rows))
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table
    header_bg = _rgb(theme.get("color_table_header_bg"))
    header_fg = _rgb(theme.get("color_table_header_text"))
    for r_idx, row in enumerate(rows):
        for c_idx in range(n_cols):
            cell = table.cell(r_idx, c_idx)
            cell.text = row[c_idx] if c_idx < len(row) else ""
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                if r_idx == 0 and header_fg:
                    p.font.color.rgb = header_fg
                    p.font.bold = True
            if r_idx == 0 and header_bg:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
    return height


def _render_slide(
    prs: Presentation,
    s: _Slide,
    images: list[tuple[bytes, int, int]],
    theme: dict,
    title_layout,
    section_layout,
    content_layout,
) -> None:
    body_font = theme.get("font_body")
    title_color = _rgb(theme.get("color_h1"))

    if s.kind == "title":
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = s.title or "Presentation"
        if title_color:
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = title_color
        # subtitle placeholder (idx 1) → any body paras (author/date/product)
        if len(slide.placeholders) > 1 and s.paras:
            _fill_text_frame(slide.placeholders[1].text_frame, s.paras, theme, body_font)
        _attach_notes(slide, s)
        return

    if s.kind == "section":
        slide = prs.slides.add_slide(section_layout)
        slide.shapes.title.text = s.title or ""
        if title_color:
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = title_color
        # Any content that followed the H1 before the next heading goes in body.
        if s.paras and len(slide.placeholders) > 1:
            _fill_text_frame(slide.placeholders[1].text_frame, s.paras, theme, body_font)
        _attach_notes(slide, s)
        return

    # content slide — Title Only layout + a manually laid-out body region
    slide = prs.slides.add_slide(content_layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = s.title or ""
        if title_color:
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = title_color

    margin = Inches(0.6)
    top = Inches(1.6)
    width = prs.slide_width - margin * 2
    bottom = prs.slide_height - Inches(0.4)
    cursor = top

    if s.paras:
        # Rough height estimate: ~0.3" per paragraph/line.
        est = Emu(int(Inches(0.34).emu * max(len(s.paras), 1)))
        box = slide.shapes.add_textbox(margin, cursor, width, min(est, bottom - cursor))
        _fill_text_frame(box.text_frame, s.paras, theme, body_font)
        cursor = Emu(cursor + est + Inches(0.1))

    for rows in s.tables:
        if cursor >= bottom:
            break
        h = _add_table(slide, rows, margin, cursor, width, theme)
        cursor = Emu(cursor + h + Inches(0.2))

    for idx in s.images:
        if idx >= len(images) or cursor >= bottom:
            continue
        png, w_px, h_px = images[idx]
        cursor = _place_picture(slide, BytesIO(png), w_px, h_px, margin, cursor, width, bottom)

    for path in s.files:
        if cursor >= bottom:
            break
        try:
            from PIL import Image

            with Image.open(path) as im:
                w_px, h_px = im.width, im.height
        except Exception:
            w_px, h_px = 800, 600
        cursor = _place_picture(slide, str(path), w_px, h_px, margin, cursor, width, bottom)

    _attach_notes(slide, s)


def _place_picture(slide, source, w_px, h_px, left, top, max_width, bottom) -> Emu:
    avail_h = bottom - top
    disp_w = min(Emu(w_px * _EMU_PER_PX), max_width)
    disp_h = Emu(int(disp_w * h_px / max(w_px, 1)))
    if disp_h > avail_h:
        disp_h = Emu(int(avail_h))
        disp_w = Emu(int(disp_h * w_px / max(h_px, 1)))
    x = Emu(int(left + (max_width - disp_w) / 2))  # centre horizontally
    slide.shapes.add_picture(source, x, top, width=disp_w, height=disp_h)
    return Emu(top + disp_h + Inches(0.15))


def _attach_notes(slide, s: _Slide) -> None:
    if s.notes:
        slide.notes_slide.notes_text_frame.text = "\n".join(s.notes)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def _strip_frontmatter(md: str) -> str:
    return re.sub(r"^---\s*\n.*?\n---\s*\n", "", md, count=1, flags=re.DOTALL)


def _extract_title(md: str) -> str | None:
    m = re.search(r"^#\s+(.+)$", md, re.MULTILINE)
    return m.group(1).strip() if m else None


def build(
    rendered_md: str,
    config: dict[str, Any],
    out_path: Path,
    *,
    doc_path: Path | None = None,
    repo_root: Path | None = None,
) -> None:
    """Convert rendered Markdown to a .pptx slide deck."""
    out_path = Path(out_path).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)

    split = str(config.get("slide_split", "h2")).lower()
    if split not in ("h1", "h2", "marker"):
        split = "h2"

    body = _strip_frontmatter(rendered_md)
    title = config.get("title") or _extract_title(body) or out_path.stem
    author = str(config.get("author", ""))
    date = str(config.get("date", ""))
    product = str(config.get("product", ""))

    # Convert Markdown → HTML, splitting on explicit slide markers first.
    # ([[field]] markers are a Word-only mechanism; markdown leaves them as
    # literal text, which is the right behaviour for slides.)
    body = _SLIDE_BREAK_RE.sub('\n\n<hr class="md-doc-slide-break">\n\n', body)
    html = markdown.Markdown(extensions=_MD_EXTENSIONS).convert(body)

    # Theme (colours + fonts from the CSS cascade) + mermaid diagram theme.
    theme = resolve_docx_theme(doc_path, repo_root) if (doc_path and repo_root) else {}
    theme = theme or {}
    mermaid_theme = None
    try:
        from .pdf import _resolve_css
        from ..mermaid import extract_theme_from_css

        css = _resolve_css(config, repo_root, doc_path=doc_path)
        if css and css.exists():
            mermaid_theme = extract_theme_from_css(css.read_text(encoding="utf-8"))
    except Exception:
        mermaid_theme = None
    html, images = _render_mermaid_to_images(html, mermaid_theme)

    # Presentation base + slide size.
    template = (
        _resolve_asset(str(config["pptx_template"]), doc_path, repo_root)
        if config.get("pptx_template")
        else None
    )
    prs = Presentation(str(template)) if template else Presentation()
    size = str(config.get("slide_size", "16:9"))
    if size in _SLIDE_SIZES and not template:
        prs.slide_width, prs.slide_height = _SLIDE_SIZES[size]

    # Parse HTML into slides; the slide-break <hr> becomes a sentinel comment
    # that the parser turns into a fresh slide boundary.
    html = html.replace('<hr class="md-doc-slide-break">', "<!--slidebreak-->")
    parser = _SlideParser(split, doc_path, repo_root)
    parser.feed(html)
    slides = [s for s in parser.slides if not s.is_empty()]

    # Prepend a title slide from config metadata.
    title_slide = _Slide(title=title, kind="title")
    meta = _build_title_meta(product, author, date)
    if meta:
        title_slide.paras = meta
    # If the body's first slide is a section made from the leading H1 that equals
    # the title, fold it into the title slide instead of duplicating it.
    if slides and slides[0].kind == "section" and (slides[0].title or "").strip() == title.strip():
        title_slide.paras = title_slide.paras or slides[0].paras
        slides = slides[1:]

    all_slides = [title_slide, *slides]

    title_layout = _layout(prs, "Title Slide", 0)
    section_layout = _layout(prs, "Section Header", 2)
    content_layout = _layout(prs, "Title Only", 5)

    for s in all_slides:
        _render_slide(prs, s, images, theme, title_layout, section_layout, content_layout)

    prs.save(str(out_path))


def _build_title_meta(product: str, author: str, date: str) -> list[_Para]:
    paras: list[_Para] = []
    if product:
        paras.append(_Para(runs=[_Run(text=product)]))
    if author:
        paras.append(_Para(runs=[_Run(text=author)]))
    if date:
        paras.append(_Para(runs=[_Run(text=date)]))
    return paras
